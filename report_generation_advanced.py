"""
Advanced Report Generation with PPT and HTML Support
Includes: PowerPoint generation, HTML reports, Executive summaries, Customizable templates
"""

import streamlit as st
import pandas as pd
import numpy as np
from datetime import datetime
import io
import base64
from pipeline_history import PipelineHistory

# Try to import report generation libraries
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from jinja2 import Template
    HAS_JINJA2 = True
except ImportError:
    HAS_JINJA2 = False

try:
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib import colors
    HAS_REPORTLAB = True
except ImportError:
    HAS_REPORTLAB = False


class AdvancedReportGenerator:
    def __init__(self):
        self.history = PipelineHistory()
    
    def render_advanced_report_ui(self):
        """Render advanced report generation interface"""
        st.subheader("📊 Advanced Report Generation")
        
        # Report format selection
        report_formats = []
        if HAS_REPORTLAB:
            report_formats.append("PDF")
        if HAS_PPTX:
            report_formats.append("PowerPoint (PPTX)")
        if HAS_JINJA2:
            report_formats.append("HTML")
        
        if not report_formats:
            st.error("No report generation libraries available.")
            st.code("pip install python-pptx jinja2 reportlab")
            return
        
        selected_format = st.selectbox("Report Format", report_formats, key="adv_report_format")
        
        # Report configuration
        col1, col2 = st.columns(2)
        
        with col1:
            st.markdown("### 📋 Report Sections")
            sections = {
                'executive_summary': st.checkbox("Executive Summary", value=True, key="rep_exec"),
                'data_overview': st.checkbox("Data Overview", value=True, key="rep_data"),
                'preprocessing': st.checkbox("Preprocessing Steps", value=True, key="rep_prep"),
                'visualizations': st.checkbox("Visualizations", value=True, key="rep_viz"),
                'model_results': st.checkbox("Model Results", value=True, key="rep_model"),
                'insights': st.checkbox("Key Insights & Recommendations", value=True, key="rep_insights")
            }
        
        with col2:
            st.markdown("### ⚙️ Report Settings")
            report_title = st.text_input("Report Title", "Data Analysis Report", key="adv_rep_title")
            author_name = st.text_input("Author", "Data Science Team", key="adv_rep_author")
            company_name = st.text_input("Company/Organization", "Your Organization", key="adv_rep_company")
            
            # Dataset selection
            if st.session_state.datasets:
                dataset_name = st.selectbox("Primary Dataset", list(st.session_state.datasets.keys()), key="adv_rep_dataset")
            else:
                dataset_name = None
                st.warning("No datasets available")
        
        # Template selection
        if selected_format == "PowerPoint (PPTX)":
            template_style = st.selectbox(
                "Presentation Template",
                ["Professional Blue", "Modern Dark", "Clean White", "Corporate"],
                key="ppt_template"
            )
        elif selected_format == "HTML":
            template_style = st.selectbox(
                "HTML Template",
                ["Bootstrap Modern", "Minimal", "Dashboard Style"],
                key="html_template"
            )
        else:
            template_style = "Standard"
        
        # Generate button
        if st.button(f"📄 Generate {selected_format} Report", type="primary", key="gen_adv_report"):
            if dataset_name:
                self._generate_advanced_report(
                    selected_format, sections, report_title, author_name,
                    company_name, dataset_name, template_style
                )
            else:
                st.error("Please select a dataset")

    
    def _generate_advanced_report(self, format_type, sections, title, author, company, dataset_name, template_style):
        """Generate advanced report in specified format"""
        try:
            with st.spinner(f"Generating {format_type} report..."):
                if format_type == "PowerPoint (PPTX)" and HAS_PPTX:
                    buffer = self._generate_pptx_report(sections, title, author, company, dataset_name, template_style)
                    file_ext = "pptx"
                    mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                
                elif format_type == "HTML" and HAS_JINJA2:
                    buffer = self._generate_html_report(sections, title, author, company, dataset_name, template_style)
                    file_ext = "html"
                    mime_type = "text/html"
                
                elif format_type == "PDF" and HAS_REPORTLAB:
                    buffer = self._generate_pdf_report(sections, title, author, company, dataset_name)
                    file_ext = "pdf"
                    mime_type = "application/pdf"
                
                else:
                    st.error(f"Format {format_type} not available")
                    return
                
                if buffer:
                    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                    filename = f"report_{timestamp}.{file_ext}"
                    
                    st.download_button(
                        label=f"📥 Download {format_type} Report",
                        data=buffer.getvalue(),
                        file_name=filename,
                        mime=mime_type
                    )
                    
                    st.success(f"✅ {format_type} report generated successfully!")
                    
                    # Log operation
                    self.history.log_step(
                        "Advanced Report Generation",
                        f"Generated {format_type} report: {title}",
                        {
                            "format": format_type,
                            "sections": list(sections.keys()),
                            "dataset": dataset_name,
                            "template": template_style
                        },
                        "success"
                    )
        
        except Exception as e:
            st.error(f"Error generating report: {str(e)}")
            self.history.log_step(
                "Advanced Report Generation",
                f"Failed to generate {format_type} report",
                {"error": str(e)},
                "error"
            )

    
    def _generate_pptx_report(self, sections, title, author, company, dataset_name, template_style):
        """Generate PowerPoint presentation"""
        try:
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Define color schemes based on template
            color_schemes = {
                "Professional Blue": {
                    'primary': RGBColor(0, 51, 102),
                    'secondary': RGBColor(0, 102, 204),
                    'accent': RGBColor(255, 153, 0)
                },
                "Modern Dark": {
                    'primary': RGBColor(33, 33, 33),
                    'secondary': RGBColor(66, 66, 66),
                    'accent': RGBColor(255, 87, 34)
                },
                "Clean White": {
                    'primary': RGBColor(51, 51, 51),
                    'secondary': RGBColor(102, 102, 102),
                    'accent': RGBColor(76, 175, 80)
                },
                "Corporate": {
                    'primary': RGBColor(0, 32, 96),
                    'secondary': RGBColor(0, 64, 128),
                    'accent': RGBColor(255, 193, 7)
                }
            }
            
            colors = color_schemes.get(template_style, color_schemes["Professional Blue"])
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title_shape = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title_shape.text = title
            subtitle.text = f"{company}\n{author}\n{datetime.now().strftime('%B %d, %Y')}"
            
            # Get dataset
            df = st.session_state.datasets[dataset_name]
            
            # Executive Summary
            if sections.get('executive_summary'):
                self._add_ppt_executive_summary(prs, df, dataset_name, colors)
            
            # Data Overview
            if sections.get('data_overview'):
                self._add_ppt_data_overview(prs, df, dataset_name, colors)
            
            # Preprocessing
            if sections.get('preprocessing'):
                self._add_ppt_preprocessing(prs, colors)
            
            # Model Results
            if sections.get('model_results'):
                self._add_ppt_model_results(prs, colors)
            
            # Key Insights
            if sections.get('insights'):
                self._add_ppt_insights(prs, df, colors)
            
            # Save to buffer
            buffer = io.BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            return buffer
        
        except Exception as e:
            st.error(f"Error generating PowerPoint: {str(e)}")
            return None

    
    def _add_ppt_executive_summary(self, prs, df, dataset_name, colors):
        """Add executive summary slide"""
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Executive Summary"
        
        # Add text box with summary
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        
        # Generate summary
        summary_points = [
            f"Dataset: {dataset_name}",
            f"Total Records: {len(df):,}",
            f"Total Features: {len(df.columns)}",
            f"Data Quality: {((df.count().sum() / df.size) * 100):.1f}% complete",
            f"Analysis Date: {datetime.now().strftime('%B %d, %Y')}"
        ]
        
        for point in summary_points:
            p = tf.add_paragraph()
            p.text = f"• {point}"
            p.level = 0
            p.font.size = Pt(16)
    
    def _add_ppt_data_overview(self, prs, df, dataset_name, colors):
        """Add data overview slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Data Overview"
        
        # Add statistics table
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        
        # Create table data
        table_data = [
            ["Metric", "Value"],
            ["Total Rows", f"{len(df):,}"],
            ["Total Columns", f"{len(df.columns)}"],
            ["Missing Values", f"{df.isnull().sum().sum():,}"],
            ["Duplicate Rows", f"{df.duplicated().sum():,}"],
            ["Memory Usage", f"{df.memory_usage(deep=True).sum() / (1024**2):.2f} MB"]
        ]
        
        table = slide.shapes.add_table(len(table_data), 2, left, top, width, height).table
        
        # Populate table
        for i, row_data in enumerate(table_data):
            for j, cell_data in enumerate(row_data):
                cell = table.cell(i, j)
                cell.text = str(cell_data)
                
                # Header row formatting
                if i == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = colors['primary']
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = RGBColor(255, 255, 255)
                            run.font.bold = True
    
    def _add_ppt_preprocessing(self, prs, colors):
        """Add preprocessing steps slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Preprocessing Steps"
        
        # Get preprocessing history
        prep_steps = [step for step in st.session_state.pipeline_history 
                     if 'preprocess' in step['operation'].lower() or 
                        'clean' in step['operation'].lower() or
                        'outlier' in step['operation'].lower()]
        
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        
        if prep_steps:
            for i, step in enumerate(prep_steps[-10:], 1):  # Last 10 steps
                p = tf.add_paragraph()
                p.text = f"{i}. {step['operation']}: {step['description'][:80]}"
                p.level = 0
                p.font.size = Pt(14)
        else:
            p = tf.add_paragraph()
            p.text = "No preprocessing steps recorded"
            p.font.size = Pt(14)
    
    def _add_ppt_model_results(self, prs, colors):
        """Add model results slide"""
        if 'training_results' not in st.session_state or not st.session_state.training_results:
            return
        
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Model Results"
        
        results = st.session_state.training_results
        
        # Create results table
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        
        table_data = [["Model", "Task", "Primary Metric", "Score"]]
        
        for result in results[-5:]:  # Last 5 models
            metrics = result['metrics']
            if result['task_type'] == 'Classification':
                score = f"{metrics.get('accuracy', 0):.4f}"
                metric_name = "Accuracy"
            else:
                score = f"{metrics.get('r2', 0):.4f}"
                metric_name = "R² Score"
            
            table_data.append([
                result['model_name'],
                result['task_type'],
                metric_name,
                score
            ])
        
        table = slide.shapes.add_table(len(table_data), 4, left, top, width, height).table
        
        for i, row_data in enumerate(table_data):
            for j, cell_data in enumerate(row_data):
                cell = table.cell(i, j)
                cell.text = str(cell_data)
                
                if i == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = colors['primary']
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = RGBColor(255, 255, 255)
                            run.font.bold = True
    
    def _add_ppt_insights(self, prs, df, colors):
        """Add key insights slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Key Insights & Recommendations"
        
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        
        # Generate insights
        insights = []
        
        # Data quality insight
        completeness = (df.count().sum() / df.size) * 100
        if completeness < 90:
            insights.append(f"Data completeness is {completeness:.1f}% - consider improving data collection")
        else:
            insights.append(f"Data quality is excellent with {completeness:.1f}% completeness")
        
        # Feature insights
        numeric_cols = df.select_dtypes(include=[np.number]).columns
        if len(numeric_cols) > 0:
            insights.append(f"Dataset contains {len(numeric_cols)} numeric features for analysis")
        
        # Missing data insight
        missing_pct = (df.isnull().sum().sum() / df.size) * 100
        if missing_pct > 5:
            insights.append(f"Missing data ({missing_pct:.1f}%) should be addressed before modeling")
        
        # Add insights to slide
        for insight in insights:
            p = tf.add_paragraph()
            p.text = f"• {insight}"
            p.level = 0
            p.font.size = Pt(16)

    
    def _generate_html_report(self, sections, title, author, company, dataset_name, template_style):
        """Generate HTML report"""
        try:
            df = st.session_state.datasets[dataset_name]
            
            # HTML template
            html_template = """
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{{ title }}</title>
    <link href="https://cdn.jsdelivr.net/npm/bootstrap@5.1.3/dist/css/bootstrap.min.css" rel="stylesheet">
    <style>
        body { font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif; }
        .header { background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); color: white; padding: 40px 0; }
        .section { margin: 30px 0; padding: 20px; background: #f8f9fa; border-radius: 10px; }
        .metric-card { background: white; padding: 20px; border-radius: 8px; box-shadow: 0 2px 4px rgba(0,0,0,0.1); margin: 10px 0; }
        .insight { background: #e3f2fd; padding: 15px; border-left: 4px solid #2196f3; margin: 10px 0; }
        table { width: 100%; margin: 20px 0; }
        th { background: #667eea; color: white; }
    </style>
</head>
<body>
    <div class="header text-center">
        <h1>{{ title }}</h1>
        <p>{{ company }} | {{ author }} | {{ date }}</p>
    </div>
    
    <div class="container mt-4">
        {% if executive_summary %}
        <div class="section">
            <h2>📊 Executive Summary</h2>
            <div class="row">
                <div class="col-md-3">
                    <div class="metric-card">
                        <h5>Total Records</h5>
                        <h3>{{ total_rows }}</h3>
                    </div>
                </div>
                <div class="col-md-3">
                    <div class="metric-card">
                        <h5>Total Features</h5>
                        <h3>{{ total_cols }}</h3>
                    </div>
                </div>
                <div class="col-md-3">
                    <div class="metric-card">
                        <h5>Data Quality</h5>
                        <h3>{{ data_quality }}%</h3>
                    </div>
                </div>
                <div class="col-md-3">
                    <div class="metric-card">
                        <h5>Missing Values</h5>
                        <h3>{{ missing_values }}</h3>
                    </div>
                </div>
            </div>
        </div>
        {% endif %}
        
        {% if data_overview %}
        <div class="section">
            <h2>📋 Data Overview</h2>
            <table class="table table-striped">
                <thead>
                    <tr>
                        <th>Metric</th>
                        <th>Value</th>
                    </tr>
                </thead>
                <tbody>
                    <tr><td>Dataset Name</td><td>{{ dataset_name }}</td></tr>
                    <tr><td>Total Rows</td><td>{{ total_rows }}</td></tr>
                    <tr><td>Total Columns</td><td>{{ total_cols }}</td></tr>
                    <tr><td>Missing Values</td><td>{{ missing_values }}</td></tr>
                    <tr><td>Duplicate Rows</td><td>{{ duplicate_rows }}</td></tr>
                    <tr><td>Memory Usage</td><td>{{ memory_usage }} MB</td></tr>
                </tbody>
            </table>
        </div>
        {% endif %}
        
        {% if preprocessing %}
        <div class="section">
            <h2>🔧 Preprocessing Steps</h2>
            {% for step in preprocessing_steps %}
            <div class="insight">
                <strong>{{ step.operation }}</strong>: {{ step.description }}
            </div>
            {% endfor %}
        </div>
        {% endif %}
        
        {% if insights %}
        <div class="section">
            <h2>💡 Key Insights & Recommendations</h2>
            {% for insight in key_insights %}
            <div class="insight">
                {{ insight }}
            </div>
            {% endfor %}
        </div>
        {% endif %}
    </div>
    
    <script src="https://cdn.jsdelivr.net/npm/bootstrap@5.1.3/dist/js/bootstrap.bundle.min.js"></script>
</body>
</html>
            """
            
            # Prepare data for template
            template_data = {
                'title': title,
                'author': author,
                'company': company,
                'date': datetime.now().strftime('%B %d, %Y'),
                'dataset_name': dataset_name,
                'total_rows': f"{len(df):,}",
                'total_cols': len(df.columns),
                'data_quality': f"{((df.count().sum() / df.size) * 100):.1f}",
                'missing_values': f"{df.isnull().sum().sum():,}",
                'duplicate_rows': f"{df.duplicated().sum():,}",
                'memory_usage': f"{df.memory_usage(deep=True).sum() / (1024**2):.2f}",
                'executive_summary': sections.get('executive_summary'),
                'data_overview': sections.get('data_overview'),
                'preprocessing': sections.get('preprocessing'),
                'insights': sections.get('insights'),
                'preprocessing_steps': st.session_state.pipeline_history[-10:] if hasattr(st.session_state, 'pipeline_history') else [],
                'key_insights': self._generate_insights(df)
            }
            
            # Render template
            template = Template(html_template)
            html_content = template.render(**template_data)
            
            # Save to buffer
            buffer = io.BytesIO()
            buffer.write(html_content.encode('utf-8'))
            buffer.seek(0)
            return buffer
        
        except Exception as e:
            st.error(f"Error generating HTML report: {str(e)}")
            return None
    
    def _generate_insights(self, df):
        """Generate key insights from data"""
        insights = []
        
        # Data quality
        completeness = (df.count().sum() / df.size) * 100
        insights.append(f"Data completeness: {completeness:.1f}%")
        
        # Missing data
        missing_pct = (df.isnull().sum().sum() / df.size) * 100
        if missing_pct > 5:
            insights.append(f"⚠️ {missing_pct:.1f}% missing data detected - consider imputation strategies")
        
        # Duplicates
        dup_count = df.duplicated().sum()
        if dup_count > 0:
            insights.append(f"⚠️ {dup_count} duplicate rows found - consider deduplication")
        
        # Feature types
        numeric_cols = len(df.select_dtypes(include=[np.number]).columns)
        categorical_cols = len(df.select_dtypes(include=['object', 'category']).columns)
        insights.append(f"Feature distribution: {numeric_cols} numeric, {categorical_cols} categorical")
        
        return insights
    
    def _generate_pdf_report(self, sections, title, author, company, dataset_name):
        """Generate PDF report (simplified version)"""
        # This would use the existing report_generation.py logic
        # For now, return a placeholder
        st.info("PDF generation uses existing report_generation.py functionality")
        return None
